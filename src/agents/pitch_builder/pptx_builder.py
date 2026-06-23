"""
src/agents/pitch_builder/pptx_builder.py
===========================================
Generador del pitchbook .pptx para el Pitch Builder Agent.

6 diapositivas, 16:9, estilo Bankinter:
    - Naranja #FF6600 + blanco #FFFFFF + grises suaves
    - Tipografia Arial / Helvetica / sans-serif

Uso:
        from src.agents.pitch_builder.pptx_builder import build_pitch_deck
        build_pitch_deck(output_path, company_name, sector, pitch, meeting_brief, model_output)
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paleta Bankinter
# ---------------------------------------------------------------------------
PRIMARY     = RGBColor(0xFF, 0x66, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY = RGBColor(0x22, 0x22, 0x22)
TEXT_SECONDARY = RGBColor(0x5A, 0x6B, 0x7B)
LIGHT_GREY  = RGBColor(0xF5, 0xF5, 0xF5)
DIVIDER     = RGBColor(0xE6, 0xE6, 0xE6)
FOOTER_GREY = RGBColor(0x8A, 0x96, 0xA3)

FONT_DISPLAY = "Arial"
FONT_BODY    = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CONTENT_X = Inches(0.78)
CONTENT_W = Inches(11.77)
HEADER_Y = Inches(0.28)
HEADER_LINE_Y = Inches(0.88)
LOGO_PATH = Path(__file__).resolve().parents[3] / "static" / "img" / "narangrande.png"


# ---------------------------------------------------------------------------
# Helpers de bajo nivel
# ---------------------------------------------------------------------------

def _add_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def _add_line(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_slide_header(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), CONTENT_X, HEADER_Y, height=Inches(0.42))
    _add_line(slide, CONTENT_X, HEADER_LINE_Y, CONTENT_W, Pt(1.5), DIVIDER)


def _add_text(slide, x, y, w, h, text, size, color, bold=False, italic=False,
               font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
               line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return tb


def _add_bullets(slide, x, y, w, h, items, size, color, font=FONT_BODY,
                  space_after=8, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return tb


def _add_numbered_list(slide, x, y, w, h, items, number_color=PRIMARY, text_color=TEXT_PRIMARY,
                       size=15, number_size=15, row_gap=0.62):
    for i, item in enumerate(items):
        row_y = y + Inches(i * row_gap)
        _add_text(slide, x, row_y, Inches(0.3), Inches(0.3),
                  f"{i + 1}.", number_size, number_color, bold=True, font=FONT_BODY)
        _add_text(slide, x + Inches(0.38), row_y, w - Inches(0.38), Inches(0.5),
                  item, size, text_color, font=FONT_BODY, line_spacing=1.2)


def _add_tag(slide, x, y, w, h, text, color=PRIMARY, size=11, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_BODY
    _add_line(slide, x, y + h - Pt(1.5), w, Pt(1.5), color)
    return shape


def _add_card(slide, x, y, w, h, fill_color=WHITE, line_color=None, shadow=True, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or DIVIDER
    shape.line.width = Pt(0.8)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _section_header(slide, number, label):
    del number
    _add_text(slide, CONTENT_X, Inches(1.1), Inches(6.5), Inches(0.38),
              label.upper(), 20, PRIMARY, bold=True, font=FONT_BODY)
    _add_line(slide, CONTENT_X, Inches(1.48), Inches(2.0), Pt(1.5), PRIMARY)


def _footer(slide, dark=False):
    del dark
    _add_text(slide, CONTENT_X, Inches(7.0), Inches(9), Inches(0.25),
              "Corporate & Deal Intelligence  \u00b7  Documento de uso interno",
              9.5, FOOTER_GREY, font=FONT_BODY)


def _safe(value, fallback=""):
    if value is None:
        return fallback
    if isinstance(value, str) and not value.strip():
        return fallback
    return value


# ---------------------------------------------------------------------------
# Constructor principal
# ---------------------------------------------------------------------------

def build_pitch_deck(
    output_path: Path,
    company_name: str,
    sector: str,
    pitch: dict,
    meeting_brief: dict | None = None,
    model_output: dict | None = None,
) -> Path:
    """
    Genera el pitchbook .pptx y lo guarda en output_path. Devuelve la ruta.
    """
    meeting_brief = meeting_brief or {}
    model_output  = model_output  or {}
    indicadores   = model_output.get("indicadores_clave") or []
    tendencias    = model_output.get("tendencias_estructurales") or []
    interpretacion = _safe(
        model_output.get("interpretacion"),
        "El equipo comercial completar\u00e1 la interpretaci\u00f3n financiera "
        "durante la preparaci\u00f3n de la reuni\u00f3n.",
    )

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    titulo          = _safe(pitch.get("titulo"), f"Propuesta comercial \u2014 {company_name}")
    subtitulo       = _safe(pitch.get("subtitulo"), "")
    oportunidad     = _safe(pitch.get("oportunidad_detectada"),
                            "Oportunidad identificada en este sector. "
                            "Los detalles se completar\u00e1n durante la preparaci\u00f3n de la reuni\u00f3n.")
    contexto_fin    = pitch.get("contexto_financiero") or []
    encaje_prod     = pitch.get("encaje_productos")    or []
    comparables     = pitch.get("comparables")         or []
    argumentos_val  = pitch.get("argumentos_valor")    or []
    proximos_pasos  = pitch.get("proximos_pasos")      or ["Agendar reuni\u00f3n de seguimiento con el cliente."]

    # =====================================================================
    # SLIDE 1 -- Portada
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _add_bg(s, WHITE)
    _add_slide_header(s)

    if sector:
        _add_tag(s, CONTENT_X, Inches(1.35), Inches(2.2), Inches(0.28),
                 sector.upper(), color=PRIMARY, size=11)

    title_size = 40 if len(titulo) <= 70 else 36
    _add_text(s, CONTENT_X, Inches(2.25), Inches(9.5), Inches(1.5),
              titulo, title_size, TEXT_PRIMARY, bold=True, font=FONT_DISPLAY, line_spacing=1.08)

    if subtitulo:
        _add_text(s, CONTENT_X, Inches(4.45), Inches(10.2), Inches(0.9),
                  subtitulo, 17, TEXT_SECONDARY, font=FONT_BODY, line_spacing=1.25)

    _footer(s)

    # =====================================================================
    # SLIDE 2 -- Oportunidad detectada
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _add_bg(s, WHITE)
    _add_slide_header(s)
    _section_header(s, 1, "Oportunidad detectada")

    _add_text(s, CONTENT_X, Inches(1.78), Inches(7.0), Inches(0.65),
              company_name, 30, TEXT_PRIMARY, bold=True, font=FONT_DISPLAY)
    _add_text(s, CONTENT_X, Inches(2.45), Inches(7.0), Inches(2.4),
              oportunidad, 15, TEXT_PRIMARY, font=FONT_BODY, line_spacing=1.35)

    if proximos_pasos:
        _add_text(s, CONTENT_X, Inches(5.75), Inches(7.0), Inches(0.3),
                  "PRÓXIMA ACCIÓN", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)
        _add_text(s, CONTENT_X, Inches(6.08), Inches(7.0), Inches(0.85),
                  proximos_pasos[0], 14, PRIMARY, font=FONT_BODY, line_spacing=1.25)

    _add_card(s, Inches(8.55), Inches(1.78), Inches(4.0), Inches(4.9), fill_color=WHITE, line_color=DIVIDER)
    _add_text(s, Inches(8.9), Inches(2.05), Inches(3.2), Inches(0.3),
              "EN UNA CIFRA", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)

    if indicadores:
        first = indicadores[0]
        val = str(first.get("valor", ""))
        val_size = 38 if len(val) <= 8 else 24
        _add_text(s, Inches(8.9), Inches(2.45), Inches(3.2), Inches(0.9),
                  val, val_size, PRIMARY, bold=True, font=FONT_DISPLAY, line_spacing=1.0)
        _add_text(s, Inches(8.9), Inches(3.35), Inches(3.2), Inches(0.7),
                  first.get("nombre", ""), 13, TEXT_SECONDARY, font=FONT_BODY, line_spacing=1.2)
    else:
        _add_text(s, Inches(8.9), Inches(2.45), Inches(3.2), Inches(0.9),
                  "—", 34, PRIMARY, bold=True, font=FONT_DISPLAY)
        _add_text(s, Inches(8.9), Inches(3.35), Inches(3.2), Inches(0.7),
                  "Sin indicador disponible", 13, TEXT_SECONDARY, font=FONT_BODY)

    _add_line(s, Inches(8.9), Inches(4.2), Inches(3.2), Pt(1), DIVIDER)
    _add_text(s, Inches(8.9), Inches(4.38), Inches(3.2), Inches(0.25),
              "SECTOR", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)
    _add_text(s, Inches(8.9), Inches(4.68), Inches(3.2), Inches(0.4),
              _safe(sector, "No especificado"), 16, TEXT_PRIMARY, bold=True, font=FONT_DISPLAY)
    _add_text(s, Inches(8.9), Inches(5.28), Inches(3.2), Inches(0.25),
              "PRIORIDAD", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)
    _add_text(s, Inches(8.9), Inches(5.58), Inches(3.2), Inches(0.35),
              "ALTA", 14, PRIMARY, bold=True, font=FONT_BODY)

    _footer(s)

    # =====================================================================
    # SLIDE 3 -- Contexto financiero
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _add_bg(s, WHITE)
    _add_slide_header(s)
    _section_header(s, 2, "Contexto financiero")

    _add_text(s, CONTENT_X, Inches(1.78), CONTENT_W, Inches(0.6),
              "Lo que dicen los números", 30, TEXT_PRIMARY, bold=True, font=FONT_DISPLAY)

    card_w, card_h, gap = Inches(3.85), Inches(2.3), Inches(0.3)
    if indicadores:
        for i, ind in enumerate(indicadores[:3]):
            x = Inches(0.7) + i * (card_w + gap)
            _add_card(s, x, Inches(2.45), card_w, card_h, fill_color=LIGHT_GREY, line_color=DIVIDER)
            value = str(ind.get("valor", ""))
            val_size = 22 if len(value) > 10 else 32
            _add_text(s, x + Inches(0.28), Inches(2.72), card_w - Inches(0.56), Inches(0.75),
                      value, val_size, PRIMARY, bold=True, font=FONT_DISPLAY, line_spacing=1.0)
            _add_text(s, x + Inches(0.28), Inches(3.95), card_w - Inches(0.56), Inches(0.6),
                      ind.get("nombre", ""), 13, TEXT_SECONDARY, font=FONT_BODY, line_spacing=1.15)
    elif contexto_fin:
        _add_card(s, CONTENT_X, Inches(2.45), CONTENT_W, card_h, fill_color=LIGHT_GREY, line_color=DIVIDER)
        _add_bullets(s, Inches(1.08), Inches(2.75), Inches(11.1), Inches(1.8),
                     contexto_fin[:4], 14, TEXT_PRIMARY, space_after=10)
    else:
        _add_card(s, CONTENT_X, Inches(2.45), CONTENT_W, card_h, fill_color=LIGHT_GREY, line_color=DIVIDER)
        _add_text(s, Inches(1.08), Inches(3.28), Inches(11.1), Inches(0.6),
                  "No se han recibido indicadores financieros estructurados para esta oportunidad.",
                  14, TEXT_SECONDARY, font=FONT_BODY)

    if tendencias:
        _add_text(s, CONTENT_X, Inches(5.25), Inches(5.7), Inches(0.3),
                  "TENDENCIAS ESTRUCTURALES", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)
        _add_bullets(s, CONTENT_X, Inches(5.58), Inches(5.7), Inches(1.35),
                     tendencias[:4], 13, TEXT_PRIMARY, space_after=8)
    elif not indicadores and contexto_fin:
        _add_text(s, CONTENT_X, Inches(5.25), Inches(5.7), Inches(0.3),
                  "RESUMEN", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)
        _add_text(s, CONTENT_X, Inches(5.58), Inches(5.7), Inches(1.15),
                  "Datos de contexto recopilados por el Earnings Reviewer Agent (ver tarjeta superior).",
                  13, TEXT_SECONDARY, font=FONT_BODY, line_spacing=1.3)

    _add_card(s, Inches(6.85), Inches(5.25), Inches(5.7), Inches(1.8), fill_color=WHITE, line_color=DIVIDER)
    _add_text(s, Inches(7.15), Inches(5.45), Inches(5.1), Inches(0.25),
              "INTERPRETACIÓN", 11, PRIMARY, bold=True, font=FONT_BODY)
    _add_text(s, Inches(7.15), Inches(5.78), Inches(5.1), Inches(1.0),
              interpretacion, 13, TEXT_PRIMARY, font=FONT_BODY, line_spacing=1.25)

    _footer(s)

    # =====================================================================
    # SLIDE 4 -- Briefing de cliente
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _add_bg(s, WHITE)
    _add_slide_header(s)
    _section_header(s, 3, "Briefing de cliente")

    _add_text(s, CONTENT_X, Inches(1.78), CONTENT_W, Inches(0.6),
              "Perfil y situación actual", 30, TEXT_PRIMARY, bold=True, font=FONT_DISPLAY)

    perfil    = _safe(meeting_brief.get("perfil"),
                       f"{company_name} es una de las oportunidades priorizadas para este sector.")
    situacion = _safe(meeting_brief.get("situacion_actual"), "")

    _add_text(s, CONTENT_X, Inches(2.35), CONTENT_W, Inches(0.9),
              perfil, 14, TEXT_PRIMARY, font=FONT_BODY, line_spacing=1.3)
    if situacion:
        _add_text(s, CONTENT_X, Inches(3.1), CONTENT_W, Inches(0.8),
                  situacion, 14, TEXT_SECONDARY, font=FONT_BODY, line_spacing=1.3)

    necesidades = meeting_brief.get("necesidades_potenciales") or ["Por definir con el equipo comercial."]
    riesgos     = meeting_brief.get("riesgos") or ["Sin riesgos identificados con la informaci\u00f3n disponible."]

    _add_card(s, Inches(0.7), Inches(4.05), Inches(5.85), Inches(2.95), fill_color=WHITE, line_color=DIVIDER)
    _add_text(s, Inches(1.0), Inches(4.35), Inches(4.9), Inches(0.3),
              "NECESIDADES POTENCIALES", 12, PRIMARY, bold=True, font=FONT_BODY)
    _add_bullets(s, Inches(1.0), Inches(4.82), Inches(5.2), Inches(1.9),
                 necesidades[:4], 13, TEXT_PRIMARY, space_after=10)

    _add_card(s, Inches(6.8), Inches(4.05), Inches(5.85), Inches(2.95), fill_color=WHITE, line_color=DIVIDER)
    _add_text(s, Inches(7.1), Inches(4.35), Inches(4.9), Inches(0.3),
              "RIESGOS Y PUNTOS DE ATENCIÓN", 12, PRIMARY, bold=True, font=FONT_BODY)
    _add_bullets(s, Inches(7.1), Inches(4.82), Inches(5.2), Inches(1.9),
                 riesgos[:4], 13, TEXT_PRIMARY, space_after=10)

    _footer(s)

    # =====================================================================
    # SLIDE 5 -- Encaje, comparables y argumentos de valor
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _add_bg(s, WHITE)
    _add_slide_header(s)
    _section_header(s, 4, "Encaje y comparables")

    _add_text(s, CONTENT_X, Inches(1.78), CONTENT_W, Inches(0.6),
              "Cómo podemos ayudar", 30, TEXT_PRIMARY, bold=True, font=FONT_DISPLAY)

    if encaje_prod:
        _add_card(s, Inches(0.7), Inches(2.4), Inches(7.0), Inches(3.15), fill_color=WHITE, line_color=DIVIDER)
        _add_text(s, Inches(1.0), Inches(2.68), Inches(6.4), Inches(0.25),
                  "ENCAJE CON PRODUCTOS DEL BANCO", 11, PRIMARY, bold=True, font=FONT_BODY)
        _add_bullets(s, Inches(1.0), Inches(3.08), Inches(6.2), Inches(2.1),
                     encaje_prod[:4], 13.5, TEXT_PRIMARY, space_after=10)

    if comparables:
        _add_card(s, Inches(8.1), Inches(2.4), Inches(4.55), Inches(3.5), fill_color=LIGHT_GREY, line_color=DIVIDER)
        _add_text(s, Inches(8.42), Inches(2.68), Inches(3.9), Inches(0.25),
                  "COMPARABLES SECTORIALES", 11, PRIMARY, bold=True, font=FONT_BODY)
        yc = Inches(3.25)
        for comp in comparables[:2]:
            _add_text(s, Inches(8.42), yc, Inches(3.9), Inches(1.5),
                      comp, 12.5, TEXT_PRIMARY, font=FONT_BODY, line_spacing=1.3)
            yc += Inches(1.4)

    if argumentos_val:
        _add_text(s, CONTENT_X, Inches(6.05), CONTENT_W, Inches(0.25),
                  "ARGUMENTOS DE VALOR", 11, TEXT_SECONDARY, bold=True, font=FONT_BODY)
        arg_w, arg_h, gap = Inches(3.85), Inches(1.0), Inches(0.3)
        for i, arg in enumerate(argumentos_val[:3]):
            x = Inches(0.7) + i * (arg_w + gap)
            _add_card(s, x, Inches(6.35), arg_w, arg_h, fill_color=WHITE, line_color=DIVIDER)
            _add_text(s, x + Inches(0.22), Inches(6.5), arg_w - Inches(0.44), Inches(0.62),
                      arg, 11.5, TEXT_PRIMARY, font=FONT_BODY, line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s)

    # =====================================================================
    # SLIDE 6 -- Proximos pasos (cierre)
    # =====================================================================
    s = prs.slides.add_slide(blank)
    _add_bg(s, WHITE)
    _add_slide_header(s)

    _add_text(s, CONTENT_X, Inches(1.55), CONTENT_W, Inches(0.7),
              "Next steps", 34, PRIMARY, bold=True, font=FONT_DISPLAY)
    _add_text(s, CONTENT_X, Inches(2.18), Inches(10.7), Inches(0.7),
              "El equipo comercial revisa, ajusta y personaliza esta propuesta "
              "antes de cualquier interacci\u00f3n con el cliente.",
              15, TEXT_SECONDARY, font=FONT_BODY, line_spacing=1.3)

    _add_numbered_list(s, CONTENT_X, Inches(3.0), Inches(10.9), Inches(3.0),
                       proximos_pasos[:4], number_color=PRIMARY, text_color=TEXT_PRIMARY,
                       size=17, number_size=17, row_gap=0.92)

    _footer(s)

    # =====================================================================
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path